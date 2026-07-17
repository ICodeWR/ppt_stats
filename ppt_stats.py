#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 统计工具
统计 PowerPoint 文件的页数、幻灯片尺寸、形状数量等信息。

依赖安装:
    pip install python-pptx
"""

import os
import sys
import json
import argparse
from pathlib import Path
from collections import defaultdict

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
except ImportError:
    print("错误: 缺少 python-pptx 依赖。请运行: pip install python-pptx")
    sys.exit(1)


def emu_to_cm(emu: int) -> float:
    """将 EMU 转换为厘米。"""
    return round(emu / 914400 * 2.54, 2)


def emu_to_px(emu: int, dpi: int = 96) -> int:
    """将 EMU 转换为像素。"""
    return int(emu / 914400 * dpi)


def get_slide_size(prs) -> dict:
    """获取幻灯片尺寸信息。"""
    width = prs.slide_width
    height = prs.slide_height
    return {
        "width_cm": emu_to_cm(width),
        "height_cm": emu_to_cm(height),
        "width_px": emu_to_px(width),
        "height_px": emu_to_px(height),
        "ratio": round(width / height, 3) if height else 0,
    }


def get_shape_info(slide) -> dict:
    """统计幻灯片中的形状信息。"""
    total = 0
    text_boxes = 0
    pictures = 0
    tables = 0
    charts = 0
    groups = 0
    others = 0
    unknown = 0

    def count_shape(shape):
        nonlocal text_boxes, pictures, tables, charts, groups, others, unknown
        try:
            if shape.has_text_frame and shape.text.strip():
                text_boxes += 1
            shape_type = shape.shape_type
            if shape_type == 13:
                pictures += 1
            elif shape_type == 19:
                tables += 1
            elif shape_type == 3:
                charts += 1
            elif shape_type == 6:
                groups += 1
            else:
                others += 1
        except (ValueError, Exception):
            unknown += 1

    try:
        for shape in slide.shapes:
            total += 1
            count_shape(shape)
    except Exception:
        pass

    return {
        "total": total,
        "text_boxes": text_boxes,
        "pictures": pictures,
        "tables": tables,
        "charts": charts,
        "groups": groups,
        "others": others,
        "unknown": unknown,
    }


def analyze_ppt(path: str) -> dict:
    """分析单个 PPT 文件。"""
    prs = Presentation(path)
    size_info = get_slide_size(prs)

    slides_info = []
    total_shapes = defaultdict(int)

    for idx, slide in enumerate(prs.slides, 1):
        info = get_shape_info(slide)
        slides_info.append({
            "slide_no": idx,
            "shapes": info,
        })
        for k, v in info.items():
            total_shapes[k] += v

    return {
        "file": os.path.basename(path),
        "file_path": path,
        "slides_count": len(prs.slides),
        "slide_width_cm": size_info["width_cm"],
        "slide_height_cm": size_info["height_cm"],
        "slide_width_px": size_info["width_px"],
        "slide_height_px": size_info["height_px"],
        "aspect_ratio": size_info["ratio"],
        "total_shapes": dict(total_shapes),
        "slides_detail": slides_info,
    }


def find_ppt_files(directory: str, recursive: bool = False) -> list:
    """查找目录下的 PPT 文件。"""
    path = Path(directory)
    pattern = "**/*.pptx" if recursive else "*.pptx"
    return [str(p) for p in path.glob(pattern) if p.is_file()]


def print_table(results: list):
    """以表格形式打印结果。"""
    if not results:
        print("未找到 PPT 文件。")
        return

    headers = ["文件名", "页数", "宽(cm)", "高(cm)", "宽(px)", "高(px)", "形状总数", "图片数", "表格数", "图表数"]
    col_widths = [max(len(h), 18) for h in headers]

    def fmt(val, width):
        return str(val).ljust(width)

    header_line = " | ".join(fmt(h, w) for h, w in zip(headers, col_widths))
    print(header_line)
    print("-" * len(header_line))

    total_slides = 0
    for r in results:
        ts = r.get("total_shapes", {})
        row = [
            r["file"],
            r["slides_count"],
            r["slide_width_cm"],
            r["slide_height_cm"],
            r["slide_width_px"],
            r["slide_height_px"],
            ts.get("total", 0),
            ts.get("pictures", 0),
            ts.get("tables", 0),
            ts.get("charts", 0),
        ]
        print(" | ".join(fmt(v, w) for v, w in zip(row, col_widths)))
        total_slides += r["slides_count"]

    print("-" * len(header_line))
    print(f"文件总数: {len(results)},  幻灯片总页数: {total_slides}")


def print_markdown_table(results: list) -> str:
    """以 Markdown 表格格式输出。"""
    if not results:
        return "未找到 PPT 文件。"

    headers = ["文件名", "页数", "宽(cm)", "高(cm)", "宽(px)", "高(px)", "形状总数", "图片数", "表格数", "图表数"]
    lines = []
    lines.append("| " + " | ".join(headers) + " |")
    lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

    for r in results:
        ts = r.get("total_shapes", {})
        row = [
            r["file"],
            str(r["slides_count"]),
            str(r["slide_width_cm"]),
            str(r["slide_height_cm"]),
            str(r["slide_width_px"]),
            str(r["slide_height_px"]),
            str(ts.get("total", 0)),
            str(ts.get("pictures", 0)),
            str(ts.get("tables", 0)),
            str(ts.get("charts", 0)),
        ]
        lines.append("| " + " | ".join(row) + " |")

    total_slides = sum(r["slides_count"] for r in results)
    lines.append(f"\n文件总数: {len(results)}, 幻灯片总页数: {total_slides}")
    return "\n".join(lines)


def print_csv(results: list) -> str:
    """以 CSV 格式输出。"""
    if not results:
        return "未找到 PPT 文件。"

    headers = ["文件名", "页数", "宽(cm)", "高(cm)", "宽(px)", "高(px)", "形状总数", "图片数", "表格数", "图表数"]
    lines = []
    lines.append(",".join(headers))

    for r in results:
        ts = r.get("total_shapes", {})
        row = [
            f'"{r["file"]}"',
            str(r["slides_count"]),
            str(r["slide_width_cm"]),
            str(r["slide_height_cm"]),
            str(r["slide_width_px"]),
            str(r["slide_height_px"]),
            str(ts.get("total", 0)),
            str(ts.get("pictures", 0)),
            str(ts.get("tables", 0)),
            str(ts.get("charts", 0)),
        ]
        lines.append(",".join(row))

    return "\n".join(lines)


def print_detail(result: dict):
    """打印单个文件的详细信息。"""
    print(f"\n文件: {result['file_path']}")
    print(f"  页数: {result['slides_count']}")
    print(f"  幻灯片尺寸: {result['slide_width_cm']}cm x {result['slide_height_cm']}cm "
          f"({result['slide_width_px']}px x {result['slide_height_px']}px)")
    print(f"  宽高比: {result['aspect_ratio']}")
    ts = result.get("total_shapes", {})
    print(f"  形状总数: {ts.get('total', 0)}")
    print(f"    - 含文本: {ts.get('text_boxes', 0)}")
    print(f"    - 图片: {ts.get('pictures', 0)}")
    print(f"    - 表格: {ts.get('tables', 0)}")
    print(f"    - 图表: {ts.get('charts', 0)}")
    print(f"    - 组合: {ts.get('groups', 0)}")
    print(f"    - 其他: {ts.get('others', 0)}")
    print(f"    - 未知: {ts.get('unknown', 0)}")

    if result.get("slides_detail"):
        print("  各页详情:")
        for s in result["slides_detail"]:
            sh = s["shapes"]
            print(f"    第 {s['slide_no']:>3} 页: 形状 {sh['total']:>3} "
                  f"(文本 {sh['text_boxes']:>3}, 图片 {sh['pictures']:>3}, "
                  f"表格 {sh['tables']:>2}, 图表 {sh['charts']:>2})")


def main():
    parser = argparse.ArgumentParser(
        description="统计 PPT 文件的页数、尺寸、形状数量等信息。"
    )
    parser.add_argument(
        "path",
        nargs="?",
        default=".",
        help="PPT 文件或目录路径 (默认: 当前目录)",
    )
    parser.add_argument(
        "-r", "--recursive",
        action="store_true",
        help="递归搜索子目录",
    )
    parser.add_argument(
        "-d", "--detail",
        action="store_true",
        help="显示每页详细信息",
    )
    parser.add_argument(
        "-j", "--json",
        action="store_true",
        help="以 JSON 格式输出",
    )
    parser.add_argument(
        "-m", "--markdown",
        action="store_true",
        help="以 Markdown 表格格式输出",
    )
    parser.add_argument(
        "-c", "--csv",
        action="store_true",
        help="以 CSV 格式输出",
    )
    parser.add_argument(
        "-o", "--output",
        metavar="FILE",
        help="将结果保存到文件",
    )
    args = parser.parse_args()

    target = args.path
    if not os.path.exists(target):
        print(f"错误: 路径不存在: {target}")
        sys.exit(1)

    files = []
    if os.path.isfile(target):
        if target.lower().endswith(".pptx"):
            files = [target]
        else:
            print(f"错误: 不是支持的 PPT 文件: {target}")
            sys.exit(1)
    else:
        files = find_ppt_files(target, recursive=args.recursive)

    if not files:
        print("未找到 .pptx 文件。")
        sys.exit(0)

    results = []
    for f in files:
        try:
            results.append(analyze_ppt(f))
        except Exception as e:
            print(f"分析失败 [{f}]: {e}")

    if args.json:
        output = json.dumps(results, ensure_ascii=False, indent=2)
        if args.output:
            with open(args.output, "w", encoding="utf-8") as f:
                f.write(output)
            print(f"结果已保存到: {args.output}")
        else:
            print(output)
    elif args.markdown:
        output = print_markdown_table(results)
        if args.output:
            with open(args.output, "w", encoding="utf-8") as f:
                f.write(output)
            print(f"结果已保存到: {args.output}")
        else:
            print(output)
    elif args.csv:
        output = print_csv(results)
        if args.output:
            with open(args.output, "w", encoding="utf-8-sig") as f:
                f.write(output)
            print(f"结果已保存到: {args.output}")
        else:
            print(output)
    else:
        if args.detail and len(results) == 1:
            print_detail(results[0])
        else:
            print_table(results)
            if args.detail:
                for r in results:
                    print_detail(r)

        if args.output:
            with open(args.output, "w", encoding="utf-8") as f:
                json.dump(results, f, ensure_ascii=False, indent=2)
            print(f"\n结果已保存到: {args.output}")


if __name__ == "__main__":
    main()
